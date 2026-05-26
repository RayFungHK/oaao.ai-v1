"""
Structured plain-text extraction + coarse segments for vault RAG ingest.

Segments become chunk boundaries inside {@code vault_document_embed} (slide / sheet / page / markdown heading /
doc flow) with metadata attached to Qdrant payloads.
"""

from __future__ import annotations

import logging
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


@dataclass
class TextSegment:
    """One logical slice of a source file (page, worksheet, slide section, markdown section, …)."""

    scope: str
    """Machine scope: pdf_page | md_section | docx_flow | docx_table | xlsx_sheet | pptx_slide | plain"""

    label: str
    """Short human / UI label."""

    body: str
    """UTF-8 text to chunk + embed."""

    meta: dict[str, Any] = field(default_factory=dict)
    """JSON-serializable Qdrant payload extras (page, slide, sheet, …)."""


def _sanitize_label(s: str, *, maxlen: int = 240) -> str:
    x = " ".join((s or "").split())
    return x[:maxlen] if len(x) > maxlen else x


def chunk_plain_text(raw: str, *, size: int, overlap: int) -> list[str]:
    """Sliding-window character chunks; avoids splitting inside math delimiters."""

    text = raw.replace("\x00", " ").strip()
    if text == "":
        return []
    protected = _math_protected_ranges(text)
    if protected:
        return _chunk_text_avoiding_ranges(text, size=size, overlap=overlap, protected=protected)
    return _chunk_text_plain(text, size=size, overlap=overlap)


_MATH_FENCE_RE = re.compile(
    r"\$\$[\s\S]*?\$\$|\\\[[\s\S]*?\\\]|"
    r"(?<!\$)\$(?:\\.|[^$\\])+\$(?!\$)|"
    r"\\\((?:\\.|[^\\])+\\\)",
)


def _math_protected_ranges(text: str) -> list[tuple[int, int]]:
    return [(m.start(), m.end()) for m in _MATH_FENCE_RE.finditer(text)]


def _inside_protected(idx: int, protected: list[tuple[int, int]]) -> bool:
    for start, end in protected:
        if start <= idx < end:
            return True
    return False


def _chunk_text_plain(text: str, *, size: int, overlap: int) -> list[str]:
    chunks: list[str] = []
    i = 0
    n = len(text)
    while i < n:
        end = min(i + size, n)
        piece = text[i:end].strip()
        if piece:
            chunks.append(piece)
        if end >= n:
            break
        step = max(1, end - overlap)
        if step <= i:
            step = end
        i = step
    return chunks


def _chunk_text_avoiding_ranges(
    text: str,
    *,
    size: int,
    overlap: int,
    protected: list[tuple[int, int]],
) -> list[str]:
    chunks: list[str] = []
    i = 0
    n = len(text)
    while i < n:
        end = min(i + size, n)
        while end < n and _inside_protected(end, protected):
            end += 1
        if end < n:
            split = text.rfind("\n\n", i + max(1, size // 3), end)
            if split > i and not _inside_protected(split, protected):
                end = split
            else:
                split = text.rfind("\n", i + max(1, size // 4), end)
                if split > i and not _inside_protected(split, protected):
                    end = split
        piece = text[i:end].strip()
        if piece:
            chunks.append(piece)
        if end >= n:
            break
        step = max(1, end - overlap)
        if step <= i:
            step = end
        i = step
    return chunks


def _format_timestamp_hms(begin_ms: int) -> str:
    total_sec = max(0, begin_ms // 1000)
    h = total_sec // 3600
    m = (total_sec % 3600) // 60
    s = total_sec % 60
    if h > 0:
        return f"{h}:{m:02d}:{s:02d}"
    return f"{m}:{s:02d}"


def build_asr_segment_pieces(
    segments: list[dict[str, Any]],
    *,
    chunk_size: int,
    max_chunks: int = 500,
) -> list[tuple[str, dict[str, Any]]]:
    """
    One ASR diarization segment per logical unit — sub-chunk only when utterance exceeds ``chunk_size``.

    Qdrant payloads include ``speaker_id``, ``begin_ms``, ``end_ms``, ``segment_index`` for transcript seek.
    """
    out: list[tuple[str, dict[str, Any]]] = []
    for seg_idx, raw in enumerate(segments):
        if not isinstance(raw, dict):
            continue
        text = str(raw.get("text") or "").replace("\x00", " ").strip()
        if not text:
            continue
        speaker_id = max(0, int(raw.get("speaker_id") or 0))
        begin_ms = max(0, int(raw.get("begin_ms") or 0))
        end_ms = max(begin_ms, int(raw.get("end_ms") or begin_ms + 500))
        speaker_label = _sanitize_label(str(raw.get("speaker_label") or f"Speaker {speaker_id + 1}"), maxlen=128)
        stamp = _format_timestamp_hms(begin_ms)
        header = f"[{speaker_label} @ {stamp}]"
        to_chunk = f"{header}\n{text}"
        sub_idx = 0
        for piece in chunk_plain_text(to_chunk, size=chunk_size, overlap=0):
            ps = piece.strip()
            if not ps:
                continue
            if len(out) >= max_chunks:
                return out
            meta: dict[str, Any] = {
                "segment_scope": "asr_transcript",
                "segment_label": speaker_label,
                "segment_index": seg_idx,
                "subchunk_index": sub_idx,
                "speaker_id": speaker_id,
                "speaker_label": speaker_label,
                "begin_ms": begin_ms,
                "end_ms": end_ms,
            }
            out.append((ps, meta))
            sub_idx += 1
    return out


def _semantic_merge_segments(segments: list[TextSegment], *, min_chars: int = 480) -> list[TextSegment]:
    """Merge adjacent small markdown sections to reduce over-fragmentation before chunking."""
    if not segments:
        return segments
    merged: list[TextSegment] = []
    buf: TextSegment | None = None
    for seg in segments:
        if seg.scope != "md_section":
            if buf:
                merged.append(buf)
                buf = None
            merged.append(seg)
            continue
        if buf and len(buf.body) + len(seg.body) < min_chars:
            buf = TextSegment(
                scope=buf.scope,
                label=buf.label,
                body=f"{buf.body.rstrip()}\n\n{seg.body.lstrip()}",
                meta=dict(buf.meta),
            )
        else:
            if buf:
                merged.append(buf)
            buf = seg
    if buf:
        merged.append(buf)
    return merged


def build_embedding_pieces(
    segments: list[TextSegment],
    *,
    chunk_size: int,
    overlap: int,
    max_chunks: int = 500,
) -> list[tuple[str, dict[str, Any]]]:
    """
    Chunk within each segment, attach structured metadata.

    Optionally prefixes each segment's blob with ``[scope: label]`` so embeddings retain slide/sheet/page context.
    """
    segments = _semantic_merge_segments(segments)
    out: list[tuple[str, dict[str, Any]]] = []
    for seg in segments:
        body = seg.body.replace("\x00", " ").strip()
        if not body:
            continue
        to_chunk = body
        if seg.scope != "plain" and (seg.label or "").strip():
            to_chunk = f"[{seg.scope}: {seg.label}]\n{body}"
        for piece in chunk_plain_text(to_chunk, size=chunk_size, overlap=overlap):
            ps = piece.strip()
            if not ps:
                continue
            if len(out) >= max_chunks:
                return out
            meta: dict[str, Any] = {
                "segment_scope": seg.scope,
                "segment_label": (seg.label or "")[:512],
            }
            for k, v in seg.meta.items():
                if v is not None and v != "":
                    meta[str(k)] = v
            out.append((ps, meta))
    return out


def resolve_document_format(path: Path, mime_type: str) -> str:
    """
    Choose extractor family. ``mime_type`` may be wrong (``application/octet-stream``); suffix wins.
    """
    suf = path.suffix.lower()
    m = (mime_type or "").lower().strip()

    if suf == ".pdf" or m == "application/pdf":
        return "pdf"
    if suf in (".md", ".markdown") or "markdown" in m:
        return "markdown"
    if suf == ".docx" or "wordprocessingml.document" in m:
        return "docx"
    if suf == ".xlsx" or "spreadsheetml.sheet" in m:
        return "xlsx"
    if suf == ".pptx" or "presentationml.presentation" in m:
        return "pptx"
    if m.startswith("text/") or m in ("application/json",) or m.endswith("+json"):
        return "text"
    if suf in {".txt", ".csv", ".log", ".json"}:
        return "text"
    return "unknown"


def extract_text_segments(path: Path, mime_type: str) -> list[TextSegment] | None:
    """
    Return ordered segments or ``None`` if the path is not a supported office / PDF / text format for this module.
    """
    kind = resolve_document_format(path, mime_type)
    if kind == "pdf":
        return _extract_pdf_segments(path)
    if kind == "markdown":
        return _extract_markdown_segments(path)
    if kind == "docx":
        return _extract_docx_segments(path)
    if kind == "xlsx":
        return _extract_xlsx_segments(path)
    if kind == "pptx":
        return _extract_pptx_segments(path)
    if kind == "text":
        return _extract_plain_segments(path)
    return None


def _extract_plain_segments(path: Path) -> list[TextSegment] | None:
    try:
        raw = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            raw = path.read_text(encoding="latin-1", errors="replace")
        except Exception:  # noqa: BLE001
            return None
    except Exception:  # noqa: BLE001
        return None
    t = raw.replace("\x00", " ").strip()
    if not t:
        return None
    return [TextSegment(scope="plain", label=path.name or "document", body=t, meta={})]


def _pypdf_page_texts(path: Path) -> list[str] | None:
    """Per-page text via pypdf; ``None`` when the reader fails (e.g. AES without cryptography)."""
    try:
        from pypdf import PdfReader  # type: ignore[import-untyped]
    except ImportError:
        logger.warning("vault_document_extract: pypdf not installed")
        return None
    try:
        reader = PdfReader(str(path))
        out: list[str] = []
        for page in reader.pages:
            tx = page.extract_text()
            body = tx.replace("\x00", " ").strip() if isinstance(tx, str) else ""
            out.append(body)
        return out
    except Exception as e:  # noqa: BLE001
        logger.warning("vault_document_extract: pypdf failed — %s (falling back to fitz/OCR)", e)
        return None


def _fitz_page_plain_text(page: Any) -> str:
    try:
        raw = page.get_text("text")
    except Exception:  # noqa: BLE001
        return ""
    return " ".join((raw or "").replace("\x00", " ").split())


def _extract_pdf_segments(path: Path) -> list[TextSegment] | None:
    min_chars = _pdf_ocr_min_chars()
    lang = _pdf_ocr_lang()
    dpi = _pdf_ocr_dpi()
    pypdf_pages = _pypdf_page_texts(path)

    fitz_doc: Any | None = None
    try:
        import fitz  # type: ignore[import-untyped]  # pymupdf

        fitz_doc = fitz.open(str(path))
    except ImportError:
        logger.info("vault_document_extract: pymupdf not installed — PDF OCR/text fallback unavailable")
    except Exception as e:  # noqa: BLE001
        logger.warning("vault_document_extract: fitz open failed — %s", e)

    if pypdf_pages is None and fitz_doc is None:
        return None

    page_count = fitz_doc.page_count if fitz_doc is not None else len(pypdf_pages or [])
    if page_count < 1:
        if fitz_doc is not None:
            fitz_doc.close()
        return None

    out: list[TextSegment] = []
    try:
        for i in range(page_count):
            body = ""
            if pypdf_pages is not None and i < len(pypdf_pages):
                body = pypdf_pages[i] or ""

            used_ocr = False
            if len(body) < min_chars and fitz_doc is not None:
                page = fitz_doc.load_page(i)
                fitz_body = _fitz_page_plain_text(page)
                if len(fitz_body) >= min_chars:
                    body = fitz_body
                else:
                    ocr_body = _pdf_page_text_ocr_from_page(page, lang=lang, dpi=dpi)
                    if ocr_body:
                        body = ocr_body
                        used_ocr = True

            if not body:
                continue
            meta: dict[str, Any] = {"page": i + 1}
            if used_ocr:
                meta["ocr"] = True
            out.append(
                TextSegment(
                    scope="pdf_page",
                    label=_sanitize_label(f"PDF page {i + 1}"),
                    body=body,
                    meta=meta,
                ),
            )
    finally:
        if fitz_doc is not None:
            fitz_doc.close()

    return out or None


def _pdf_ocr_min_chars() -> int:
    try:
        return max(8, min(500, int(os.environ.get("OAAO_VAULT_PDF_OCR_MIN_CHARS", "48") or "48")))
    except ValueError:
        return 48


def _pdf_ocr_lang() -> str:
    return (os.environ.get("OAAO_VAULT_PDF_OCR_LANG") or "eng+chi_tra").strip() or "eng"


def _pdf_ocr_dpi() -> int:
    try:
        return max(120, min(400, int(os.environ.get("OAAO_VAULT_PDF_OCR_DPI", "200") or "200")))
    except ValueError:
        return 200


def _pdf_page_text_ocr_from_page(page: Any, *, lang: str, dpi: int) -> str | None:
    """OCR one already-open fitz page (scanned PDFs)."""
    try:
        import pytesseract  # type: ignore[import-untyped]
        from PIL import Image
    except ImportError:
        logger.info(
            "vault_document_extract: PDF OCR skipped — install pymupdf, pytesseract, Pillow (+ tesseract binary)",
        )
        return None

    try:
        pix = page.get_pixmap(dpi=dpi, alpha=False)
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        text = pytesseract.image_to_string(img, lang=lang)
        cleaned = " ".join((text or "").replace("\x00", " ").split())
        return cleaned if cleaned else None
    except Exception as e:  # noqa: BLE001
        page_no = int(getattr(page, "number", 0)) + 1
        logger.warning("vault_document_extract: PDF OCR page %s failed — %s", page_no, e)
        return None


_HEADING_MD = re.compile(r"^(#{1,6})\s+(.+?)\s*$")


def _extract_markdown_segments(path: Path) -> list[TextSegment] | None:
    raw = _read_utf8_fallback(path)
    if raw is None:
        return None
    lines = raw.splitlines()
    out: list[TextSegment] = []
    bucket: list[str] = []
    label = "preamble"
    level = 0

    def flush() -> None:
        nonlocal bucket, label, level
        if not bucket:
            return
        body = "\n".join(bucket).strip()
        bucket = []
        if not body:
            return
        meta: dict[str, Any] = {}
        if level:
            meta["heading_level"] = level
        out.append(TextSegment(scope="md_section", label=_sanitize_label(label), body=body, meta=meta))

    for line in lines:
        m = _HEADING_MD.match(line.rstrip())
        if m:
            flush()
            level = len(m.group(1))
            label = m.group(2).strip()
            bucket.append(line.strip())
        else:
            bucket.append(line)
    flush()
    return out or None


def _read_utf8_fallback(path: Path) -> str | None:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            return path.read_text(encoding="latin-1", errors="replace")
        except Exception:  # noqa: BLE001
            return None
    except Exception:  # noqa: BLE001
        return None


def _extract_docx_segments(path: Path) -> list[TextSegment] | None:
    try:
        from docx import Document  # type: ignore[import-untyped]
        from docx.document import Document as DocumentType  # type: ignore[import-untyped]
        from docx.oxml.ns import qn  # type: ignore[import-untyped]
        from docx.table import Table  # type: ignore[import-untyped]
        from docx.text.paragraph import Paragraph  # type: ignore[import-untyped]
    except ImportError:
        logger.warning("vault_document_extract: python-docx not installed")
        return None

    def iter_block_items(parent: DocumentType) -> Any:
        for child in parent.element.body.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, parent)
            elif child.tag == qn("w:tbl"):
                yield Table(child, parent)

    def is_heading_para(para: Paragraph) -> bool:
        name = (getattr(para.style, "name", "") or "").lower()
        return "heading" in name or name.startswith("title")

    try:
        doc = Document(str(path))
        segments_out: list[TextSegment] = []
        buf: list[str] = []
        heading = ""

        def flush_buffer(note: str) -> None:
            nonlocal buf, heading
            if not buf:
                return
            block = "\n".join(buf).strip()
            buf = []
            if not block:
                return
            label = heading if heading else note
            meta = {"docx_heading": heading} if heading else {}
            segments_out.append(
                TextSegment(
                    scope="docx_flow",
                    label=_sanitize_label(label),
                    body=block,
                    meta=meta,
                ),
            )

        for block in iter_block_items(doc):
            if isinstance(block, Paragraph):
                txt = (block.text or "").strip()
                if not txt:
                    continue
                if is_heading_para(block):
                    flush_buffer("section")
                    heading = txt
                    buf.append(txt)
                else:
                    buf.append(txt)
            elif isinstance(block, Table):
                flush_buffer("body")
                rows_out: list[str] = []
                for row in block.rows:
                    cells = [(c.text or "").strip().replace("\n", " ") for c in row.cells]
                    if any(cells):
                        rows_out.append(" | ".join(cells))
                if rows_out:
                    segments_out.append(
                        TextSegment(
                            scope="docx_table",
                            label=_sanitize_label(f"Table after «{heading or 'document'}»"),
                            body="\n".join(rows_out),
                            meta={"after_heading": heading} if heading else {},
                        ),
                    )

        flush_buffer("document")
        return segments_out or None
    except Exception as e:  # noqa: BLE001
        logger.warning("vault_document_extract: DOCX failed — %s", e)
        return None


def _extract_xlsx_segments(path: Path) -> list[TextSegment] | None:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError:
        logger.warning("vault_document_extract: openpyxl not installed")
        return None

    try:
        max_rows = max(512, min(500_000, int((os.environ.get("OAAO_VAULT_XLSX_MAX_ROWS", "25000") or "25000"))))
        wb = load_workbook(str(path), read_only=True, data_only=True)
        out: list[TextSegment] = []
        try:
            for sheet in wb.worksheets:
                rows: list[str] = []
                for row_ix, row in enumerate(sheet.iter_rows(values_only=True)):
                    if row_ix >= max_rows:
                        logger.warning(
                            "vault_document_extract: sheet %r truncated at %s rows (OAAO_VAULT_XLSX_MAX_ROWS)",
                            sheet.title,
                            max_rows,
                        )
                        break
                    cells = []
                    for c in row:
                        if c is None:
                            cells.append("")
                        else:
                            cells.append(str(c).strip())
                    line = " | ".join(cells)
                    if line.strip("| \t"):
                        rows.append(line)
                body = "\n".join(rows).strip()
                if body:
                    name = str(sheet.title or "Sheet")
                    out.append(
                        TextSegment(
                            scope="xlsx_sheet",
                            label=_sanitize_label(f"Sheet «{name}»"),
                            body=body,
                            meta={"sheet": name},
                        ),
                    )
        finally:
            wb.close()
        return out or None
    except Exception as e:  # noqa: BLE001
        logger.warning("vault_document_extract: XLSX failed — %s", e)
        return None


def _extract_pptx_segments(path: Path) -> list[TextSegment] | None:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]
    except ImportError:
        logger.warning("vault_document_extract: python-pptx not installed")
        return None

    try:
        prs = Presentation(str(path))
        out: list[TextSegment] = []
        for i, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tbl = shape.table  # type: ignore[attr-defined]
                    for row in tbl.rows:
                        cells = [(c.text_frame.text or "").strip().replace("\n", " ") for c in row.cells]
                        if any(cells):
                            parts.append(" | ".join(cells))
                elif hasattr(shape, "text"):
                    t = str(shape.text or "").strip()
                    if t:
                        parts.append(t)
            body = "\n".join(parts).strip()
            if body:
                out.append(
                    TextSegment(
                        scope="pptx_slide",
                        label=_sanitize_label(f"Slide {i}"),
                        body=body,
                        meta={"slide": i},
                    ),
                )
        return out or None
    except Exception as e:  # noqa: BLE001
        logger.warning("vault_document_extract: PPTX failed — %s", e)
        return None
