"""PPTX locale + font inventory for import analyze (deterministic, pre-LLM)."""

from __future__ import annotations

import logging
import re
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

logger = logging.getLogger(__name__)

# OOXML drawingml namespace
_DML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

_CJK_RE = re.compile(
    r"[\u4e00-\u9fff\u3400-\u4dbf\uf900-\ufaff"
    r"\u3040-\u309f\u30a0-\u30ff"
    r"\uac00-\ud7af]"
)
_LATIN_RE = re.compile(r"[A-Za-z]")
_DIGIT_RE = re.compile(r"[0-9]")

# Rough trad vs simp discriminators (not exhaustive)
_TRAD_BIAS = set("臺為裡麼兩與種國會體經說這還來個們時麼嗎臺灣軟體網路資訊")
_SIMP_BIAS = set("台为里么两个与国家会体经说这还来们时么吗台湾软件网络信息")

_LATIN_ONLY_TYPEFACES = frozenset(
    {
        "arial",
        "helvetica",
        "calibri",
        "calibri light",
        "times new roman",
        "georgia",
        "verdana",
        "tahoma",
        "trebuchet ms",
        "courier new",
        "garamond",
        "cambria",
        "century gothic",
        "franklin gothic",
        "segoe ui",
        "roboto",
        "open sans",
        "lato",
    }
)

_CJK_TYPEFACE_MARKERS = (
    "jhenghei",
    "yahei",
    "pingfang",
    "noto sans c",
    "noto sans tc",
    "noto sans sc",
    "source han",
    "mingliu",
    "pmingliu",
    "dfkai",
    "kaiti",
    "fangsong",
    "simhei",
    "simsun",
    "hiragino",
    "yu gothic",
    "malgun",
    "apple sd gothic",
)

_FONT_STACKS: dict[str, str] = {
    "zh-Hant": (
        '"Microsoft JhengHei", "PingFang TC", "Noto Sans TC", '
        '"Segoe UI", system-ui, sans-serif'
    ),
    "zh-Hans": (
        '"Microsoft YaHei", "PingFang SC", "Noto Sans SC", '
        '"Segoe UI", system-ui, sans-serif'
    ),
    "en": (
        'system-ui, -apple-system, "Segoe UI", Roboto, Helvetica, Arial, sans-serif'
    ),
    "ja": ('"Hiragino Sans", "Yu Gothic UI", "Noto Sans JP", sans-serif'),
    "ko": ('"Apple SD Gothic Neo", "Malgun Gothic", "Noto Sans KR", sans-serif'),
    "mixed": (
        '"Microsoft JhengHei", "PingFang TC", "Noto Sans TC", '
        '"Segoe UI", system-ui, sans-serif'
    ),
}


def _text_corpus_from_profile(profile: dict[str, Any]) -> str:
    parts: list[str] = []
    slides = profile.get("slides")
    if isinstance(slides, list):
        for row in slides:
            if not isinstance(row, dict):
                continue
            for key in ("text_sample", "title_guess"):
                val = row.get(key)
                if isinstance(val, str) and val.strip():
                    parts.append(val.strip())
    geom_slides = profile.get("geometry_slides")
    if isinstance(geom_slides, list):
        for row in geom_slides:
            if not isinstance(row, dict):
                continue
            for slot in row.get("geometry_slots") or []:
                if isinstance(slot, dict):
                    t = slot.get("text")
                    if isinstance(t, str) and t.strip():
                        parts.append(t.strip())
    return "\n".join(parts)


def detect_locale(text: str) -> dict[str, Any]:
    """Infer primary locale from slide text (deterministic)."""
    corpus = text or ""
    cjk = len(_CJK_RE.findall(corpus))
    latin = len(_LATIN_RE.findall(corpus))
    digits = len(_DIGIT_RE.findall(corpus))
    total = max(1, cjk + latin + digits)

    trad_hits = sum(1 for ch in corpus if ch in _TRAD_BIAS)
    simp_hits = sum(1 for ch in corpus if ch in _SIMP_BIAS)

    kana = len(re.findall(r"[\u3040-\u309f\u30a0-\u30ff]", corpus))
    hangul = len(re.findall(r"[\uac00-\ud7af]", corpus))
    han = len(re.findall(r"[\u4e00-\u9fff]", corpus))

    cjk_ratio = cjk / total
    latin_ratio = latin / total

    primary = "en"
    secondary: list[str] = []
    script_mix = "latin_primary"
    confidence = 0.5

    if hangul > max(han, latin) * 0.4 and hangul >= 8:
        primary = "ko"
        script_mix = "cjk_primary"
        confidence = min(0.95, 0.6 + hangul / total)
    elif kana > max(han, latin) * 0.25 and kana >= 6:
        primary = "ja"
        script_mix = "cjk_primary"
        confidence = min(0.95, 0.6 + kana / total)
    elif cjk_ratio >= 0.22 or han >= 12:
        if trad_hits > simp_hits * 1.2:
            primary = "zh-Hant"
        elif simp_hits > trad_hits * 1.2:
            primary = "zh-Hans"
        else:
            primary = "zh-Hant"
        script_mix = "cjk_primary"
        if latin_ratio >= 0.12:
            script_mix = "zh_latin_mixed"
            secondary.append("en")
        confidence = min(0.98, 0.55 + cjk_ratio)
    elif latin_ratio >= 0.45:
        primary = "en"
        script_mix = "latin_primary"
        if cjk_ratio >= 0.08:
            script_mix = "zh_latin_mixed"
            secondary.append("zh-Hant")
        confidence = min(0.92, 0.5 + latin_ratio)
    else:
        primary = "mixed"
        script_mix = "mixed"
        confidence = 0.45

    return {
        "primary": primary,
        "secondary": secondary,
        "confidence": round(confidence, 3),
        "script_mix": script_mix,
        "char_counts": {
            "cjk": cjk,
            "latin": latin,
            "digits_punct": digits,
            "kana": kana,
            "hangul": hangul,
        },
    }


def text_has_cjk(text: str) -> bool:
    return bool(_CJK_RE.search(text or ""))


def _is_cjk_capable_typeface(name: str) -> bool:
    low = (name or "").strip().lower()
    if not low:
        return False
    return any(marker in low for marker in _CJK_TYPEFACE_MARKERS)


def _is_latin_only_typeface(name: str) -> bool:
    low = (name or "").strip().lower()
    return low in _LATIN_ONLY_TYPEFACES or (
        low.startswith("arial") and "unicode" not in low
    )


def extract_font_typefaces(pptx_path: Path) -> dict[str, Any]:
    """Collect typeface names from runs + theme XML + font table."""
    used: set[str] = set()
    theme_major = ""
    theme_minor = ""
    has_embedded = False
    pptx_font_files: list[str] = []

    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            for name in zf.namelist():
                if name.startswith("ppt/fonts/") and name.endswith(".fntdata"):
                    has_embedded = True
                    pptx_font_files.append(Path(name).name)
            theme_paths = sorted(p for p in zf.namelist() if p.startswith("ppt/theme/theme") and p.endswith(".xml"))
            if theme_paths:
                root = ET.fromstring(zf.read(theme_paths[0]))
                for tag, key in (("majorFont", "theme_major"), ("minorFont", "theme_minor")):
                    node = root.find(f".//a:fontScheme/a:{tag}/a:latin", _DML_NS)
                    if node is not None:
                        tf = (node.get("typeface") or "").strip()
                        if tf:
                            used.add(tf)
                            if key == "theme_major":
                                theme_major = tf
                            else:
                                theme_minor = tf
                    ea = root.find(f".//a:fontScheme/a:{tag}/a:ea", _DML_NS)
                    if ea is not None:
                        tf = (ea.get("typeface") or "").strip()
                        if tf:
                            used.add(tf)
    except (OSError, zipfile.BadZipFile, ET.ParseError) as exc:
        logger.warning("pptx_theme_font_parse_failed path=%s err=%s", pptx_path, exc)

    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation(str(pptx_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                try:
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            name = (run.font.name or "").strip()
                            if name:
                                used.add(name)
                        if not para.runs:
                            name = (para.font.name or "").strip() if hasattr(para, "font") else ""
                            if name:
                                used.add(name)
                except Exception:  # noqa: BLE001
                    continue
    except ImportError:
        logger.info("python-pptx unavailable for font run scan")
    except Exception:  # noqa: BLE001
        logger.exception("pptx_font_run_scan_failed path=%s", pptx_path)

    ranked = sorted(used, key=str.lower)
    return {
        "used_typefaces": ranked,
        "theme_major": theme_major,
        "theme_minor": theme_minor,
        "has_embedded": has_embedded,
        "pptx_font_files": pptx_font_files,
    }


def _font_locale_mismatch(primary: str, fonts: dict[str, Any]) -> str | None:
    if primary not in ("zh-Hant", "zh-Hans", "ja", "ko"):
        return None
    used = fonts.get("used_typefaces")
    if not isinstance(used, list) or not used:
        return None
    has_cjk_font = any(_is_cjk_capable_typeface(str(n)) for n in used)
    if has_cjk_font:
        return None
    latin_only = all(_is_latin_only_typeface(str(n)) for n in used if str(n).strip())
    if latin_only:
        return (
            f"PPTX uses Latin-only fonts ({', '.join(str(x) for x in used[:3])}) "
            f"but slide text is primarily {primary}; HTML output will use a CJK-safe stack."
        )
    return None


def build_typography_hints(locale: dict[str, Any], fonts: dict[str, Any]) -> dict[str, Any]:
    primary = str(locale.get("primary") or "en")
    script_mix = str(locale.get("script_mix") or "latin_primary")
    stack = _FONT_STACKS.get(primary) or _FONT_STACKS["mixed"]

    used = fonts.get("used_typefaces")
    if isinstance(used, list):
        for name in used:
            n = str(name).strip()
            if _is_cjk_capable_typeface(n) and n not in stack:
                stack = f'"{n}", {stack}'
                break

    line_height = 1.45 if primary.startswith("zh") or primary in ("ja", "ko") else 1.35
    letter_spacing = "normal"
    avoid: list[str] = []
    if primary.startswith("zh"):
        avoid = [str(n) for n in (used or []) if _is_latin_only_typeface(str(n))][:6]

    hints: dict[str, Any] = {
        "recommended_stack": stack,
        "primary_locale": primary,
        "script_mix": script_mix,
        "line_height_factor": line_height,
        "letter_spacing": letter_spacing,
        "avoid_typefaces": avoid,
    }
    mismatch = _font_locale_mismatch(primary, fonts)
    if mismatch:
        hints["locale_font_mismatch"] = mismatch
    return hints


def enrich_profile_typography(pptx_path: Path, profile: dict[str, Any]) -> dict[str, Any]:
    """Attach locale, fonts, typography_hints to PPTX profile."""
    corpus = _text_corpus_from_profile(profile)
    locale = detect_locale(corpus)
    fonts = extract_font_typefaces(pptx_path)
    hints = build_typography_hints(locale, fonts)
    out = dict(profile)
    out["locale"] = locale
    out["fonts"] = fonts
    out["typography_hints"] = hints
    return out


def apply_typography_to_deck_style(
    deck_style: dict[str, Any],
    profile: dict[str, Any],
    *,
    llm_typography: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Merge LLM typography under deterministic locale/font constraints."""
    out = dict(deck_style)
    hints = profile.get("typography_hints")
    locale = profile.get("locale")
    if not isinstance(hints, dict):
        return out

    typo = dict(out.get("typography") or {})
    if isinstance(llm_typography, dict):
        typo.update({k: v for k, v in llm_typography.items() if v is not None})

    primary = str(hints.get("primary_locale") or (locale or {}).get("primary") or "en")
    typo["primary_locale"] = primary
    typo["script_mix"] = str(hints.get("script_mix") or "")
    typo["font_stack"] = str(hints.get("recommended_stack") or typo.get("font_stack") or "")
    if hints.get("line_height_factor"):
        typo["line_height_factor"] = hints["line_height_factor"]
    if hints.get("letter_spacing"):
        typo["letter_spacing"] = hints["letter_spacing"]

    mismatch = hints.get("locale_font_mismatch")
    if isinstance(mismatch, str) and mismatch.strip():
        typo["locale_font_warning"] = mismatch.strip()

    out["typography"] = typo
    principles = list(out.get("design_principles") or [])
    lang_rule = (
        f"Primary language: {primary} — use the locked font_stack; "
        "do not switch to Latin-only fonts for CJK body text."
    )
    if lang_rule not in principles:
        principles.insert(0, lang_rule)
    out["design_principles"] = principles[:8]

    slide_prompt = str(out.get("slide_prompt") or "")
    if primary.startswith("zh") and "zh" not in slide_prompt.lower():
        out["slide_prompt"] = (
            f"{slide_prompt.rstrip()} Write slide copy in Traditional Chinese "
            f"when the user uses Chinese ({primary})."
        ).strip()

    return out


def pptx_master_locale_css(deck_style: dict[str, Any] | None) -> str:
    """Extra CSS for pptx_master shells from deck typography locale."""
    typo = {}
    if isinstance(deck_style, dict):
        raw = deck_style.get("typography")
        if isinstance(raw, dict):
            typo = raw
    primary = str(typo.get("primary_locale") or "")
    lh = typo.get("line_height_factor")
    try:
        line_height = float(lh) if lh is not None else 1.35
    except (TypeError, ValueError):
        line_height = 1.35
    stack = str(typo.get("font_stack") or "").strip()
    if not stack and primary.startswith("zh"):
        stack = _FONT_STACKS.get("zh-Hant" if primary == "zh-Hant" else "zh-Hans", _FONT_STACKS["zh-Hant"])
    elif not stack and primary in ("ja", "ko"):
        stack = _FONT_STACKS.get(primary, _FONT_STACKS["zh-Hant"])

    rules = []
    if stack:
        rules.append(f".oaao-layout-pptx_master {{ font-family: {stack}; }}")
        rules.append(f".oaao-pptx-slot-inner {{ font-family: {stack}; }}")
    if primary.startswith("zh") or primary in ("ja", "ko"):
        rules.append(
            f".oaao-pptx-slot-inner {{ line-height: {line_height}; "
            "word-break: break-word; overflow-wrap: anywhere; }}"
        )
        rules.append(".oaao-pptx-slot-inner { letter-spacing: normal; }")
    html_lang = primary if primary else "en"
    return "\n".join(rules), html_lang
