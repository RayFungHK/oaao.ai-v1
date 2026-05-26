"""Extract a compact structural profile from PPTX for template analysis (no rendering)."""

from __future__ import annotations

import logging
import re
from collections import Counter
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


def extract_pptx_profile(path: Path, *, max_slides: int = 12) -> dict[str, Any]:
    """
    Build JSON-safe summary: slide count, dimensions, per-slide text samples.
    Requires python-pptx (same as vault embed).
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]
    except ImportError as exc:
        raise RuntimeError("python-pptx not installed in orchestrator image") from exc

    prs = Presentation(str(path))
    w_emu = int(getattr(prs.slide_width, "emu", 0) or 0)
    h_emu = int(getattr(prs.slide_height, "emu", 0) or 0)
    slides_out: list[dict[str, Any]] = []

    for i, slide in enumerate(prs.slides, start=1):
        if i > max_slides:
            break
        parts: list[str] = []
        has_table = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                has_table = True
                tbl = shape.table  # type: ignore[attr-defined]
                for row in tbl.rows:
                    cells = [
                        (c.text_frame.text or "").strip().replace("\n", " ") for c in row.cells
                    ]
                    if any(cells):
                        parts.append(" | ".join(cells))
            elif hasattr(shape, "text"):
                t = str(shape.text or "").strip()
                if t:
                    parts.append(t)
        body = "\n".join(parts).strip()
        title_guess = _first_line(body) or f"Slide {i}"
        bullets = len(re.findall(r"^[-•*]\s", body, re.MULTILINE))
        slides_out.append(
            {
                "index": i,
                "title_guess": title_guess[:120],
                "text_sample": body[:600],
                "bullet_count": bullets,
                "has_table": has_table,
            }
        )

    palette_hints = _extract_palette_hints(prs)

    return {
        "slide_count": len(prs.slides),
        "slide_width_emu": w_emu,
        "slide_height_emu": h_emu,
        "aspect": _aspect_label(w_emu, h_emu),
        "slides": slides_out,
        "palette_hints": palette_hints,
    }


def _hex_from_rgb(rgb: Any) -> str | None:
    try:
        if rgb is None:
            return None
        return f"#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}"
    except (TypeError, ValueError, IndexError):
        return None


def _luminance(hex_color: str) -> float:
    h = (hex_color or "").lstrip("#")
    if len(h) != 6:
        return 0.5
    try:
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
    except ValueError:
        return 0.5


def _collect_slide_colors(slide: Any) -> list[str]:
    """Sample solid fills from slide background and shapes."""
    from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]

    found: list[str] = []

    def add_hex(h: str | None) -> None:
        if h and re.fullmatch(r"#[0-9a-fA-F]{6}", h):
            found.append(h.lower())

    try:
        fill = slide.background.fill
        if fill.type is not None and hasattr(fill, "fore_color"):
            fc = fill.fore_color
            if fc.type == MSO_COLOR_TYPE.RGB and fc.rgb is not None:
                add_hex(_hex_from_rgb(fc.rgb))
    except Exception:  # noqa: BLE001
        pass

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            continue
        try:
            if not hasattr(shape, "fill"):
                continue
            fill = shape.fill
            if fill.type is None or not hasattr(fill, "fore_color"):
                continue
            fc = fill.fore_color
            if fc.type == MSO_COLOR_TYPE.RGB and fc.rgb is not None:
                add_hex(_hex_from_rgb(fc.rgb))
            elif fc.type == MSO_COLOR_TYPE.SCHEME and hasattr(fc, "theme_color"):
                if fc.theme_color == MSO_THEME_COLOR.BACKGROUND_1:
                    add_hex("#ffffff")
                elif fc.theme_color == MSO_THEME_COLOR.TEXT_1:
                    add_hex("#111827")
        except Exception:  # noqa: BLE001
            continue
    return found


def _extract_palette_hints(prs: Any, *, max_slides: int = 6) -> dict[str, str]:
    """Dominant colors from imported deck — feeds template analyze + preview CSS."""
    counter: Counter[str] = Counter()
    for i, slide in enumerate(prs.slides, start=1):
        if i > max_slides:
            break
        for hex_c in _collect_slide_colors(slide):
            counter[hex_c] += 1

    if not counter:
        return {}

    ranked = [c for c, _ in counter.most_common(12)]
    by_lum = sorted(ranked, key=_luminance)
    bg = by_lum[0]
    fg = by_lum[-1] if len(by_lum) > 1 else "#0f172a"
    accent = (
        by_lum[len(by_lum) // 2]
        if len(by_lum) > 2
        else (ranked[1] if len(ranked) > 1 else "#2563eb")
    )
    if _luminance(accent) < 0.25:
        accent = by_lum[-1]
    muted = by_lum[len(by_lum) // 3] if len(by_lum) > 3 else "#64748b"
    card = f"rgba({_hex_to_rgb_parts(bg)},0.08)" if bg.startswith("#") else "rgba(15,23,42,0.06)"

    if _luminance(bg) > 0.72:
        fg, bg = bg, fg if _luminance(fg) < 0.45 else "#0f172a"

    return {
        "bg": bg,
        "fg": fg,
        "muted": muted,
        "accent": accent,
        "card": card,
        "bar": accent,
    }


def _hex_to_rgb_parts(hex_color: str) -> str:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return "15,23,42"
    return f"{int(h[0:2], 16)},{int(h[2:4], 16)},{int(h[4:6], 16)}"


def profile_slides_for_template(template: dict[str, Any]) -> list[dict[str, Any]]:
    profile = template.get("profile")
    if not isinstance(profile, dict):
        return []
    slides = profile.get("slides")
    if not isinstance(slides, list):
        return []
    return [s for s in slides if isinstance(s, dict)]


def theme_from_profile(profile: dict[str, Any]) -> dict[str, str]:
    """Merge palette_hints into a full theme dict for deck_style / template.theme."""
    hints = profile.get("palette_hints")
    if not isinstance(hints, dict):
        return {}
    out: dict[str, str] = {}
    for key in ("bg", "fg", "muted", "accent", "card", "bar"):
        val = hints.get(key)
        if isinstance(val, str) and val.strip():
            out[key] = val.strip()
    return out


def _first_line(text: str) -> str:
    for line in (text or "").split("\n"):
        s = line.strip()
        if s:
            return _strip_md_inline(s)
    return ""


def _strip_md_inline(text: str) -> str:
    s = (text or "").strip()
    s = re.sub(r"\*\*([^*]+)\*\*", r"\1", s)
    s = re.sub(r"\*([^*]+)\*", r"\1", s)
    return s.strip()


def _aspect_label(w_emu: int, h_emu: int) -> str:
    if w_emu <= 0 or h_emu <= 0:
        return "16:9"
    ratio = w_emu / h_emu
    if abs(ratio - 16 / 9) < 0.08:
        return "16:9"
    if abs(ratio - 4 / 3) < 0.08:
        return "4:3"
    return f"{w_emu}:{h_emu}"
