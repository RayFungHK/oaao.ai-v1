"""Export on-disk slide HTML decks to a real ``.pptx`` (LibreOffice + python-pptx)."""

from __future__ import annotations

import io
import json
import logging
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

from oaao_orchestrator.slide_project.pptx_render import (
    _pdftoppm_binary,
    _soffice_binary,
    pptx_render_available,
)

logger = logging.getLogger(__name__)

# 13.333" × 7.5" — 16:9 (matches 1280×720 canvas)
_SLIDE_WIDTH_EMU = 12192000
_SLIDE_HEIGHT_EMU = 6858000

_SLOT_ID_HTML_RE = re.compile(
    r'data-slot-id="([^"]+)"[^>]*>.*?<div class="oaao-pptx-slot-inner">(.*?)</div>',
    re.I | re.S,
)
_MD_BULLET_RE = re.compile(r"^[-*•]\s+", re.M)
_HTML_TAG_RE = re.compile(r"<[^>]+>")


def pptx_export_enabled() -> bool:
    raw = (os.environ.get("OAAO_PPTX_EXPORT") or "1").strip().lower()
    if raw in ("0", "false", "no", "off"):
        return False
    return pptx_render_available()


def _plain_slot_text(raw: str) -> str:
    text = _MD_BULLET_RE.sub("", (raw or "").strip())
    text = _HTML_TAG_RE.sub("", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text[:2000]


def _slot_values_for_slide(slide_dir: Path, html_path: Path | None) -> dict[str, str]:
    slots_path = slide_dir / "slots.json"
    if slots_path.is_file():
        try:
            data = json.loads(slots_path.read_text(encoding="utf-8"))
            raw = data.get("slots") if isinstance(data, dict) else None
            if isinstance(raw, dict):
                return {
                    str(k): _plain_slot_text(str(v))
                    for k, v in raw.items()
                    if str(v).strip()
                }
        except Exception:  # noqa: BLE001
            logger.warning("pptx_export_slots_json_read_failed path=%s", slots_path)
    if html_path and html_path.is_file():
        try:
            doc = html_path.read_text(encoding="utf-8")
        except OSError:
            return {}
        out: dict[str, str] = {}
        for sid, inner in _SLOT_ID_HTML_RE.findall(doc):
            plain = _plain_slot_text(inner)
            if plain:
                out[str(sid).strip()] = plain
        return out
    return {}


def _bbox_center(bbox: dict[str, float]) -> tuple[float, float]:
    left = float(bbox.get("left_pct") or 0)
    top = float(bbox.get("top_pct") or 0)
    w = float(bbox.get("width_pct") or 0)
    h = float(bbox.get("height_pct") or 0)
    return left + w / 2.0, top + h / 2.0


def _match_slot_id(
    shape_bbox: dict[str, float],
    geometry_slots: list[dict[str, Any]],
    *,
    max_dist_pct: float = 18.0,
) -> str | None:
    cx, cy = _bbox_center(shape_bbox)
    best_id: str | None = None
    best_d = max_dist_pct
    for row in geometry_slots:
        if not isinstance(row, dict):
            continue
        sid = str(row.get("slot_id") or "").strip()
        if not sid:
            continue
        gx, gy = _bbox_center(row)
        d = ((cx - gx) ** 2 + (cy - gy) ** 2) ** 0.5
        if d < best_d:
            best_d = d
            best_id = sid
    return best_id


def _set_shape_text(shape: Any, text: str) -> None:
    try:
        if not hasattr(shape, "text_frame"):
            return
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
    except Exception:  # noqa: BLE001
        try:
            shape.text = text  # type: ignore[attr-defined]
        except Exception:  # noqa: BLE001
            logger.debug("pptx_export_set_text_failed", exc_info=True)


def export_from_template_source(
    *,
    source_pptx: Path,
    project_dir: Path,
    slides_spec: list[dict[str, Any]],
    slide_count: int,
) -> bytes | None:
    """
    Copy archived template PPTX and fill text shapes using geometry + ``slots.json``.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        logger.info("pptx_export_template_skip no python-pptx")
        return None

    if not source_pptx.is_file():
        return None

    with tempfile.TemporaryDirectory(prefix="oaao-pptx-export-") as tmp:
        dest = Path(tmp) / "deck.pptx"
        shutil.copy2(source_pptx, dest)
        prs = Presentation(str(dest))
        w_emu = int(getattr(prs.slide_width, "emu", 0) or 0)
        h_emu = int(getattr(prs.slide_height, "emu", 0) or 0)
        if w_emu < 1 or h_emu < 1:
            w_emu, h_emu = _SLIDE_WIDTH_EMU, _SLIDE_HEIGHT_EMU

        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]
        from oaao_orchestrator.slide_project.pptx_geometry import _bbox_pct, _shape_text  # noqa: PLC0415

        updated = 0
        for spec in sorted(slides_spec, key=lambda s: int(s.get("index") or 0)):
            idx = int(spec.get("index") or 0)
            if idx < 1 or idx > slide_count:
                continue
            if idx > len(prs.slides):
                break
            slide = prs.slides[idx - 1]
            geom = spec.get("geometry_slots")
            if not isinstance(geom, list) or not geom:
                continue
            slide_dir = project_dir / f"slides/{idx:02d}"
            html_path = project_dir / f"slides/{idx:02d}/slide.html"
            values = _slot_values_for_slide(slide_dir, html_path)
            if not values:
                continue

            stack: list[Any] = list(slide.shapes)
            while stack:
                shp = stack.pop()
                try:
                    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                        stack.extend(shp.shapes)
                        continue
                except Exception:  # noqa: BLE001
                    continue
                text = _shape_text(shp)
                if not text:
                    continue
                bbox = _bbox_pct(shp, w_emu, h_emu)
                sid = _match_slot_id(bbox, geom)
                if not sid:
                    continue
                new_text = values.get(sid)
                if not new_text:
                    continue
                _set_shape_text(shp, new_text)
                updated += 1

        if updated < 1:
            logger.info("pptx_export_template_no_shapes_updated")
            return None

        # Trim extra slides when deck is shorter than template file.
        if slide_count < len(prs.slides):
            _trim_slides(prs, slide_count)

        buf = io.BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        logger.info("pptx_export_template_ok bytes=%s updated=%s", len(data), updated)
        return data


def _trim_slides(prs: Any, keep: int) -> None:
    """Remove slides after ``keep`` (best-effort)."""
    if keep < 1:
        return
    try:
        sld_id_lst = prs.slides._sldIdLst  # type: ignore[attr-defined]
        nodes = list(sld_id_lst)
        for node in reversed(nodes[keep:]):
            r_id = node.rId
            prs.part.drop_rel(r_id)
            sld_id_lst.remove(node)
    except Exception:  # noqa: BLE001
        logger.warning("pptx_export_trim_slides_failed keep=%s", keep, exc_info=True)


def _html_to_png(html_path: Path, out_dir: Path) -> Path | None:
    """HTML → PDF (LibreOffice) → PNG (pdftoppm) — direct PNG export fails for slide HTML."""
    soffice = _soffice_binary()
    pdftoppm = _pdftoppm_binary()
    if not soffice or not pdftoppm:
        return None
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            [
                soffice,
                "--headless",
                "--nologo",
                "--nofirststartwizard",
                "--convert-to",
                "pdf",
                "--outdir",
                str(out_dir),
                str(html_path.resolve()),
            ],
            check=False,
            capture_output=True,
            timeout=120,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        logger.warning("pptx_export_html_to_pdf_failed path=%s err=%s", html_path, exc)
        return None

    pdf_path = out_dir / f"{html_path.stem}.pdf"
    if not pdf_path.is_file():
        pdfs = sorted(out_dir.glob("*.pdf"), key=lambda p: p.stat().st_mtime, reverse=True)
        pdf_path = pdfs[0] if pdfs else None
    if pdf_path is None or not pdf_path.is_file():
        return None

    png_base = out_dir / "slide"
    try:
        subprocess.run(
            [
                pdftoppm,
                "-png",
                "-singlefile",
                "-r",
                "150",
                str(pdf_path),
                str(png_base),
            ],
            check=False,
            capture_output=True,
            timeout=60,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        logger.warning("pptx_export_pdf_to_png_failed path=%s err=%s", pdf_path, exc)
        return None

    out_png = png_base.with_suffix(".png")
    if out_png.is_file() and out_png.stat().st_size > 100:
        return out_png
    pngs = sorted(out_dir.glob("*.png"), key=lambda p: p.stat().st_mtime, reverse=True)
    return pngs[0] if pngs else None


def export_from_slide_html_images(
    *,
    project_dir: Path,
    pages: list[dict[str, Any]],
) -> bytes | None:
    """Render each ``slides/NN/slide.html`` to PNG via LibreOffice, assemble with python-pptx."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return None

    soffice = _soffice_binary()
    if not soffice:
        return None

    ordered = sorted(
        [p for p in pages if isinstance(p, dict)],
        key=lambda p: int(p.get("index") or 0),
    )
    if not ordered:
        return None

    with tempfile.TemporaryDirectory(prefix="oaao-pptx-html-") as tmp:
        tmp_path = Path(tmp)
        pngs: list[Path] = []
        for page in ordered:
            idx = int(page.get("index") or 0)
            if idx < 1:
                continue
            rel = str(page.get("html_path") or f"slides/{idx:02d}/slide.html").strip()
            html_path = project_dir / rel
            if not html_path.is_file():
                logger.warning("pptx_export_missing_html idx=%s path=%s", idx, html_path)
                continue
            slide_tmp = tmp_path / f"slide_{idx:02d}"
            slide_tmp.mkdir(parents=True, exist_ok=True)
            png = _html_to_png(html_path, slide_tmp)
            if png is not None:
                pngs.append(png)

        if not pngs:
            logger.warning("pptx_export_no_pngs")
            return None

        prs = Presentation()
        prs.slide_width = _SLIDE_WIDTH_EMU
        prs.slide_height = _SLIDE_HEIGHT_EMU
        blank = prs.slide_layouts[6]
        for png in pngs:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(png),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        buf = io.BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        logger.info("pptx_export_html_images_ok slides=%s bytes=%s", len(pngs), len(data))
        return data


def build_project_pptx(
    *,
    project_dir: Path,
    deck_title: str,
    pages: list[dict[str, Any]],
    slides_spec: list[dict[str, Any]],
    template_source_pptx: Path | None = None,
) -> bytes | None:
    """
    Build export bytes: template text fill (when shapes match), else HTML→PNG deck.
    """
    if not pptx_export_enabled():
        logger.info("pptx_export_disabled")
        return None

    slide_count = len(pages) or len(slides_spec)
    if template_source_pptx and template_source_pptx.is_file() and slides_spec:
        try:
            data = export_from_template_source(
                source_pptx=template_source_pptx,
                project_dir=project_dir,
                slides_spec=slides_spec,
                slide_count=slide_count,
            )
            if data and len(data) > 2000:
                return data
        except Exception:  # noqa: BLE001
            logger.exception("pptx_export_template_failed title=%s", deck_title)

    try:
        return export_from_slide_html_images(project_dir=project_dir, pages=pages)
    except Exception:  # noqa: BLE001
        logger.exception("pptx_export_html_images_failed title=%s", deck_title)
        return None
