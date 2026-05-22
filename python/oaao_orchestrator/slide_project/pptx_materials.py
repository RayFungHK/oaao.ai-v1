"""CP2 — unpack PPTX media + picture-shape positions into materials/manifest.json."""

from __future__ import annotations

import json
import logging
import os
import re
import shutil
import zipfile
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

MANIFEST_VERSION = 1
MANIFEST_REL = "materials/manifest.json"
MEDIA_PREFIX = "materials/media/"


def pptx_materials_enabled() -> bool:
    raw = (os.environ.get("OAAO_PPTX_MATERIALS") or "1").strip().lower()
    return raw not in ("0", "false", "no", "off")


def _overlay_raster_on_fidelity() -> bool:
    raw = (os.environ.get("OAAO_PPTX_MATERIAL_OVERLAY_RASTER") or "").strip().lower()
    return raw in ("1", "true", "yes", "on")


def _safe_media_name(name: str) -> str:
    base = Path(name).name
    base = re.sub(r"[^a-zA-Z0-9._-]", "_", base)
    return base[:120] if base else "asset.bin"


def _extract_zip_media(pptx_path: Path, media_dir: Path) -> list[str]:
    """Copy ``ppt/media/*`` into ``materials/media/``; return relative paths under materials/."""
    media_dir.mkdir(parents=True, exist_ok=True)
    written: list[str] = []
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            for info in zf.infolist():
                if not info.filename.startswith("ppt/media/") or info.is_dir():
                    continue
                safe = _safe_media_name(info.filename.split("/")[-1])
                dest = media_dir / safe
                if dest.is_file():
                    rel = f"{MEDIA_PREFIX}{safe}"
                    if rel not in written:
                        written.append(rel)
                    continue
                try:
                    data = zf.read(info.filename)
                    dest.write_bytes(data)
                    written.append(f"{MEDIA_PREFIX}{safe}")
                except OSError as exc:
                    logger.warning("pptx_material_zip_write_failed file=%s err=%s", safe, exc)
    except (OSError, zipfile.BadZipFile) as exc:
        logger.warning("pptx_material_zip_open_failed path=%s err=%s", pptx_path, exc)
    return sorted(set(written))


def _bbox_pct(shape: Any, slide_w: int, slide_h: int) -> dict[str, float]:
    if slide_w < 1 or slide_h < 1:
        return {"left_pct": 0.0, "top_pct": 0.0, "width_pct": 100.0, "height_pct": 100.0}
    left = float(shape.left) / slide_w * 100.0
    top = float(shape.top) / slide_h * 100.0
    width = float(shape.width) / slide_w * 100.0
    height = float(shape.height) / slide_h * 100.0
    return {
        "left_pct": round(max(0.0, min(100.0, left)), 2),
        "top_pct": round(max(0.0, min(100.0, top)), 2),
        "width_pct": round(max(1.0, min(100.0 - left, width)), 2),
        "height_pct": round(max(1.0, min(100.0 - top, height)), 2),
    }


def _picture_kind(ext: str) -> str:
    e = (ext or "").lower().lstrip(".")
    if e == "svg":
        return "svg"
    if e in ("png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff"):
        return "image"
    return "media"


def _iter_flat_shapes(container: Any) -> Any:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]
    except ImportError:
        yield from list(container.shapes)
        return

    stack: list[Any] = list(container.shapes)
    while stack:
        shape = stack.pop()
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                stack.extend(shape.shapes)
                continue
        except Exception:  # noqa: BLE001
            pass
        yield shape


def _blip_embed_ids(shape: Any) -> list[str]:
    try:
        from pptx.oxml.ns import qn  # type: ignore[import-untyped]
    except ImportError:
        return []

    ids: list[str] = []
    for blip in shape._element.xpath(".//*[local-name()='blip']"):
        rid = blip.get(qn("r:embed"))
        if rid and rid not in ids:
            ids.append(rid)
    return ids


def _blob_from_slide_rid(slide: Any, rid: str) -> tuple[bytes | None, str]:
    try:
        part = slide.part.related_part(rid)
    except KeyError:
        return None, "png"
    blob = getattr(part, "blob", None)
    if not isinstance(blob, (bytes, bytearray)) or not blob:
        return None, "png"
    partname = str(getattr(part, "partname", "") or "")
    ext = Path(partname).suffix.lstrip(".").lower() or "png"
    return bytes(blob), ext


def _write_media_blob(
    *,
    media_dir: Path,
    slide_index: int,
    asset_n: int,
    blob: bytes,
    ext: str,
    preferred_name: str = "",
) -> str:
    if preferred_name:
        safe = _safe_media_name(preferred_name)
    else:
        safe = _safe_media_name(f"slide{slide_index:02d}_asset{asset_n:02d}.{ext}")
    dest = media_dir / safe
    if not dest.is_file():
        dest.write_bytes(blob)
    return f"{MEDIA_PREFIX}{safe}"


def _append_visual_asset(
    assets: list[dict[str, Any]],
    *,
    slide_index: int,
    asset_n: int,
    rel_path: str,
    kind: str,
    bbox: dict[str, float],
    zip_media_paths: set[str],
    used_paths: set[str],
) -> None:
    if rel_path in used_paths:
        return
    area = float(bbox["width_pct"]) * float(bbox["height_pct"])
    if area < 0.08:
        return
    used_paths.add(rel_path)
    zip_media_paths.discard(rel_path)
    assets.append(
        {
            "id": f"asset_{slide_index}_{asset_n}",
            "kind": kind,
            "path": rel_path,
            "left_pct": bbox["left_pct"],
            "top_pct": bbox["top_pct"],
            "width_pct": bbox["width_pct"],
            "height_pct": bbox["height_pct"],
            "z_index": 2 + asset_n,
            "overlay_on_fidelity": kind == "svg" or _overlay_raster_on_fidelity(),
        }
    )


def _extract_slide_picture_assets(
    slide: Any,
    *,
    slide_index: int,
    slide_w_emu: int,
    slide_h_emu: int,
    media_dir: Path,
    zip_media_paths: set[str],
) -> list[dict[str, Any]]:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]
    except ImportError:
        return []

    assets: list[dict[str, Any]] = []
    used_paths: set[str] = set()
    asset_n = 0

    for shape in _iter_flat_shapes(slide):
        bbox = _bbox_pct(shape, slide_w_emu, slide_h_emu)
        rel_path = ""
        kind = "image"
        ext = "png"

        # Standard picture shapes
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = str(getattr(image, "ext", None) or "png").lower()
                kind = _picture_kind(ext)
                rel_path = _write_media_blob(
                    media_dir=media_dir,
                    slide_index=slide_index,
                    asset_n=asset_n + 1,
                    blob=image.blob,
                    ext=ext,
                )
        except Exception:  # noqa: BLE001
            rel_path = ""

        # Google Slides / placeholder blips (a:blip r:embed)
        if not rel_path:
            for rid in _blip_embed_ids(shape):
                blob, ext = _blob_from_slide_rid(slide, rid)
                if not blob:
                    continue
                kind = _picture_kind(ext)
                try:
                    part = slide.part.related_part(rid)
                    preferred = Path(str(getattr(part, "partname", "") or "")).name
                except KeyError:
                    preferred = ""
                rel_path = _write_media_blob(
                    media_dir=media_dir,
                    slide_index=slide_index,
                    asset_n=asset_n + 1,
                    blob=blob,
                    ext=ext,
                    preferred_name=preferred,
                )
                break

        if not rel_path:
            continue

        asset_n += 1
        _append_visual_asset(
            assets,
            slide_index=slide_index,
            asset_n=asset_n,
            rel_path=rel_path,
            kind=kind,
            bbox=bbox,
            zip_media_paths=zip_media_paths,
            used_paths=used_paths,
        )

    return assets


def build_materials_manifest(
    pptx_path: Path,
    asset_dir: Path,
    *,
    slide_count: int | None = None,
) -> dict[str, Any] | None:
    """
    Unpack ``ppt/media`` and picture shapes into ``asset_dir/materials/``.
    Returns manifest dict or None when disabled / failed.
    """
    if not pptx_materials_enabled():
        return None
    if not pptx_path.is_file():
        return None

    materials_dir = asset_dir / "materials"
    media_dir = materials_dir / "media"
    if media_dir.exists():
        shutil.rmtree(media_dir, ignore_errors=True)
    media_dir.mkdir(parents=True, exist_ok=True)

    zip_paths = set(_extract_zip_media(pptx_path, media_dir))
    slides_out: list[dict[str, Any]] = []

    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        logger.info("python-pptx unavailable for pptx_materials")
        prs = None
    else:
        prs = Presentation(str(pptx_path))

    if prs is not None:
        w_emu = int(getattr(prs.slide_width, "emu", 0) or 0)
        h_emu = int(getattr(prs.slide_height, "emu", 0) or 0)
        for i, slide in enumerate(prs.slides, start=1):
            if slide_count is not None and i > slide_count:
                break
            assets = _extract_slide_picture_assets(
                slide,
                slide_index=i,
                slide_w_emu=w_emu,
                slide_h_emu=h_emu,
                media_dir=media_dir,
                zip_media_paths=zip_paths,
            )
            slides_out.append(
                {
                    "index": i,
                    "background": {"type": "render_png", "path": f"render/{i:02d}.png"},
                    "assets": assets,
                }
            )

    orphan_assets: list[dict[str, Any]] = []
    for n, rel in enumerate(sorted(zip_paths)):
        ext = Path(rel).suffix.lower()
        kind = _picture_kind(ext.lstrip("."))
        orphan_assets.append(
            {
                "id": f"zip_{n + 1}",
                "kind": kind,
                "path": rel,
                "overlay_on_fidelity": kind == "svg" or _overlay_raster_on_fidelity(),
            }
        )

    manifest: dict[str, Any] = {
        "version": MANIFEST_VERSION,
        "media_root": "materials/media",
        "slides": slides_out,
        "orphan_media": orphan_assets,
    }

    manifest_path = asset_dir / MANIFEST_REL
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    logger.info(
        "pptx_materials_ok path=%s slides=%s orphans=%s",
        pptx_path,
        len(slides_out),
        len(orphan_assets),
    )
    return manifest


def apply_pptx_materials(
    pptx_path: Path,
    row: dict[str, Any],
    asset_dir: Path,
) -> dict[str, Any] | None:
    """Build manifest under template asset dir; set row metadata fields."""
    profile = row.get("profile")
    slide_count = None
    if isinstance(profile, dict):
        slides = profile.get("slides")
        if isinstance(slides, list):
            slide_count = len(slides)

    manifest = build_materials_manifest(pptx_path, asset_dir, slide_count=slide_count)
    if manifest is None:
        return None

    row["materials_manifest"] = MANIFEST_REL
    row["materials_version"] = MANIFEST_VERSION
    return manifest
