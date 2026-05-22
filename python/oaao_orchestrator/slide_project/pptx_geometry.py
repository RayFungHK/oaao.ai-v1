"""Phase 3 — extract positioned text shapes from PPTX (bbox % + slot roles)."""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

MAX_GEOMETRY_SLOTS = 8
MIN_BOX_AREA_PCT = 0.12


def pptx_master_enabled() -> bool:
    import os

    raw = (os.environ.get("OAAO_PPTX_MASTER") or "1").strip().lower()
    return raw not in ("0", "false", "no", "off")


def _placeholder_slot_id(shape: Any) -> str | None:
    try:
        if not shape.is_placeholder:
            return None
        ph = shape.placeholder_format
        ptype = ph.type
        name = str(ptype).lower() if ptype is not None else ""
    except Exception:  # noqa: BLE001
        return None

    if "title" in name and "subtitle" not in name and "body" not in name:
        return "title"
    if "subtitle" in name:
        return "subtitle"
    if "body" in name or "content" in name or "object" in name:
        return "body"
    if "footer" in name:
        return "footer"
    if "header" in name:
        return "header"
    return None


def _rgb_font_color_to_hex(color: Any) -> str:
    """OOXML RGB run color → ``#rrggbb``."""
    if color is None:
        return ""
    try:
        from pptx.dml.color import RGBColor  # type: ignore[import-untyped]

        if isinstance(color, RGBColor):
            return f"#{int(color[0]):02x}{int(color[1]):02x}{int(color[2]):02x}"
    except Exception:  # noqa: BLE001
        pass
    raw = str(color).strip().lstrip("#")
    if len(raw) == 6 and re.fullmatch(r"[0-9a-fA-F]{6}", raw):
        return f"#{raw.lower()}"
    return ""


def _shape_typography(shape: Any) -> dict[str, Any]:
    """Font size / weight / family / color from the dominant text run."""
    meta: dict[str, Any] = {}
    try:
        if not getattr(shape, "has_text_frame", False):
            return meta
        tf = shape.text_frame
        max_pt = 0.0
        bold = False
        family = ""
        color_hex = ""
        align = ""
        for para in tf.paragraphs:
            if not align and para.alignment is not None:
                try:
                    from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]

                    mapping = {
                        PP_ALIGN.LEFT: "left",
                        PP_ALIGN.CENTER: "center",
                        PP_ALIGN.RIGHT: "right",
                        PP_ALIGN.JUSTIFY: "justify",
                    }
                    align = mapping.get(para.alignment, "")
                except Exception:  # noqa: BLE001
                    align = ""
            runs = list(para.runs) or [para]
            for run in runs:
                font = getattr(run, "font", None)
                if font is None:
                    continue
                if font.size is not None and getattr(font.size, "pt", None):
                    max_pt = max(max_pt, float(font.size.pt))
                if font.bold:
                    bold = True
                if font.name and not family:
                    family = str(font.name).strip()
                try:
                    if font.color is not None and font.color.type is not None:
                        rgb = getattr(font.color, "rgb", None)
                        hx = _rgb_font_color_to_hex(rgb)
                        if hx:
                            color_hex = hx
                except Exception:  # noqa: BLE001
                    pass
        if max_pt > 0:
            meta["font_pt"] = round(max_pt, 1)
        if bold:
            meta["font_weight"] = 700
        if family:
            meta["font_family"] = family[:80]
        if color_hex:
            meta["color"] = color_hex
        if align:
            meta["text_align"] = align
    except Exception:  # noqa: BLE001
        return meta
    return meta


def _shape_text(shape: Any) -> str:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl = shape.table  # type: ignore[attr-defined]
            rows: list[str] = []
            for row in tbl.rows:
                cells = [(c.text_frame.text or "").strip() for c in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            return "\n".join(rows).strip()
        if hasattr(shape, "text"):
            return str(shape.text or "").strip()
    except Exception:  # noqa: BLE001
        return ""
    return ""


def _bbox_pct(shape: Any, slide_w: int, slide_h: int) -> dict[str, float]:
    if slide_w < 1 or slide_h < 1:
        return {"left_pct": 0.0, "top_pct": 0.0, "width_pct": 100.0, "height_pct": 100.0}
    left = float(shape.left) / slide_w * 100.0
    top = float(shape.top) / slide_h * 100.0
    width = float(shape.width) / slide_w * 100.0
    height = float(shape.height) / slide_h * 100.0
    return {
        "left_pct": round(max(0.0, min(100.0, left)), 2),
        "top_pct": round(max(0.0, min(100.0, top)), 2),
        "width_pct": round(max(1.0, min(100.0 - left, width)), 2),
        "height_pct": round(max(1.0, min(100.0 - top, height)), 2),
    }


def _infer_role_from_geometry(
    *,
    bbox: dict[str, float],
    text: str,
    order: int,
    total: int,
) -> str:
    top = float(bbox.get("top_pct") or 0)
    height = float(bbox.get("height_pct") or 0)
    width = float(bbox.get("width_pct") or 0)
    lines = [ln for ln in text.split("\n") if ln.strip()]
    bullet_lines = sum(1 for ln in lines if re.match(r"^[-•*]\s", ln.strip()))

    if order == 0 and top < 22 and height < 22:
        return "title"
    if top < 18 and height < 14 and len(text) < 80:
        return "subtitle" if order > 0 else "title"
    if bullet_lines >= 2 or len(lines) >= 4:
        return "bullets"
    if width < 45 and top > 25:
        return "callout"
    if height > 35:
        return "body"
    return f"slot_{order + 1}"


def _unique_slot_id(base: str, used: set[str]) -> str:
    sid = re.sub(r"[^a-z0-9_]", "", (base or "slot").lower()) or "slot"
    if sid not in used:
        used.add(sid)
        return sid
    n = 2
    while f"{sid}_{n}" in used:
        n += 1
    sid = f"{sid}_{n}"
    used.add(sid)
    return sid


def extract_slide_geometry_shapes(
    slide: Any,
    *,
    slide_w_emu: int,
    slide_h_emu: int,
) -> list[dict[str, Any]]:
    """Text/table shapes with percentage bbox and stable slot_id."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]
    except ImportError:
        return []

    candidates: list[dict[str, Any]] = []
    stack: list[Any] = list(slide.shapes)

    while stack:
        shape = stack.pop()
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                stack.extend(shape.shapes)
                continue
        except Exception:  # noqa: BLE001
            continue

        text = _shape_text(shape)
        if not text or len(text) < 2:
            continue

        bbox = _bbox_pct(shape, slide_w_emu, slide_h_emu)
        area = float(bbox["width_pct"]) * float(bbox["height_pct"])
        if area < MIN_BOX_AREA_PCT:
            continue

        ph_slot = _placeholder_slot_id(shape)
        candidates.append(
            {
                "text": text[:1200],
                "bbox": bbox,
                "area": area,
                "placeholder_slot": ph_slot,
                "shape_name": str(getattr(shape, "name", "") or "")[:80],
                "shape": shape,
            }
        )

    candidates.sort(key=lambda c: (-float(c.get("area") or 0), float(c["bbox"]["top_pct"])))
    used: set[str] = set()
    out: list[dict[str, Any]] = []
    for i, row in enumerate(candidates[:MAX_GEOMETRY_SLOTS]):
        bbox = row["bbox"]
        ph = row.get("placeholder_slot")
        role = ph or _infer_role_from_geometry(
            bbox=bbox,
            text=str(row.get("text") or ""),
            order=i,
            total=len(candidates),
        )
        slot_id = _unique_slot_id(str(role), used)
        entry: dict[str, Any] = {
            "slot_id": slot_id,
            "left_pct": bbox["left_pct"],
            "top_pct": bbox["top_pct"],
            "width_pct": bbox["width_pct"],
            "height_pct": bbox["height_pct"],
            "text": row["text"],
            "placeholder": ph or "",
            "role": ph or role,
        }
        src_shape = row.get("shape")
        if src_shape is not None:
            typo = _shape_typography(src_shape)
            if typo:
                entry.update(typo)
        out.append(entry)
    return out


def enrich_profile_with_geometry(path: Path, profile: dict[str, Any]) -> dict[str, Any]:
    """Add ``geometry_slots`` per slide in profile (in-place copy returned)."""
    if not pptx_master_enabled():
        return profile

    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return profile

    out = dict(profile)
    prs = Presentation(str(path))
    w_emu = int(getattr(prs.slide_width, "emu", 0) or 0)
    h_emu = int(getattr(prs.slide_height, "emu", 0) or 0)
    slides = out.get("slides")
    if not isinstance(slides, list):
        return out

    enriched: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        prof = slides[i - 1] if i - 1 < len(slides) else {"index": i}
        if not isinstance(prof, dict):
            prof = {"index": i}
        row = dict(prof)
        row.setdefault("index", i)
        geom = extract_slide_geometry_shapes(slide, slide_w_emu=w_emu, slide_h_emu=h_emu)
        if geom:
            row["geometry_slots"] = geom
            row["geometry_mode"] = "pptx_master"
        enriched.append(row)

    out["slides"] = enriched
    out["geometry_version"] = 1
    return out
